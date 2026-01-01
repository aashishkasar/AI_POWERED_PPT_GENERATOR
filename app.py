from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

# Create a presentation object
prs = Presentation()

# Define the white background color
white_color = RGBColor(0xFF, 0xFF, 0xFF)

# Function to apply background color to a slide
def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Slide 1: Title Slide
slide_layout_title = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(slide_layout_title)
set_slide_background(slide1, white_color)

title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "Agentic AI with Langchain"
subtitle1.text = "A Simple Chatbot Example"

# Slide 2: What is Agentic AI?
slide_layout_content = prs.slide_layouts[1] # Title and Content layout
slide2 = prs.slides.add_slide(slide_layout_content)
set_slide_background(slide2, white_color)

title2 = slide2.shapes.title
title2.text = "Understanding Agentic AI"

body2 = slide2.shapes.placeholders[1]
tf2 = body2.text_frame
tf2.clear() # Clear existing text

p2_1 = tf2.add_paragraph()
p2_1.text = "Goal-Oriented: AI that breaks down complex tasks into sub-goals."
p2_1.level = 0

p2_2 = tf2.add_paragraph()
p2_2.text = "Reasoning: Uses a Large Language Model (LLM) to decide what to do next."
p2_2.level = 0

p2_3 = tf2.add_paragraph()
p2_3.text = "Tool Use: Can select and use external tools (e.g., search, calculator) to achieve goals."
p2_3.level = 0

p2_4 = tf2.add_paragraph()
p2_4.text = "Iterative Process: Plans, acts, observes, and refines its approach."
p2_4.level = 0

# Slide 3: Langchain & Agents
slide3 = prs.slides.add_slide(slide_layout_content)
set_slide_background(slide3, white_color)

title3 = slide3.shapes.title
title3.text = "Langchain's Role in Agentic AI"

body3 = slide3.shapes.placeholders[1]
tf3 = body3.text_frame
tf3.clear()

p3_1 = tf3.add_paragraph()
p3_1.text = "Framework for developing applications powered by LLMs."
p3_1.level = 0

p3_2 = tf3.add_paragraph()
p3_2.text = "Agents: Core component enabling LLMs to interact with their environment."
p3_2.level = 0

p3_3 = tf3.add_paragraph()
p3_3.text = "Key Components: LLM (the brain), Tools (the actions), Agent Executor (orchestration)."
p3_3.level = 0

p3_4 = tf3.add_paragraph()
p3_4.text = "Simplifies dynamic decision-making and tool selection for complex tasks."
p3_4.level = 0

# Slide 4: Simple Chatbot Example
slide4 = prs.slides.add_slide(slide_layout_content)
set_slide_background(slide4, white_color)

title4 = slide4.shapes.title
title4.text = "Agentic Chatbot in Action"

body4 = slide4.shapes.placeholders[1]
tf4 = body4.text_frame
tf4.clear()

p4_1 = tf4.add_paragraph()
p4_1.text = "Goal: Answer user questions, even if requiring external knowledge."
p4_1.level = 0

p4_2 = tf4.add_paragraph()
p4_2.text = "Tools Used:"
p4_2.level = 0
p4_2_1 = tf4.add_paragraph()
p4_2_1.text = "- LLM (e.g., OpenAI GPT-3.5) for reasoning."
p4_2_1.level = 1
p4_2_2 = tf4.add_paragraph()
p4_2_2.text = "- Search Tool (e.g., Google Search API) for external data."
p4_2_2.level = 1

p4_3 = tf4.add_paragraph()
p4_3.text = "Example Process:"
p4_3.level = 0
p4_3_1 = tf4.add_paragraph()
p4_3_1.text = "1. User asks: 'What is the population of Tokyo?'"
p4_3_1.level = 1
p4_3_2 = tf4.add_paragraph()
p4_3_2.text = "2. Agent reasons: 'I need a search tool to find this information.'"
p4_3_2.level = 1
p4_3_3 = tf4.add_paragraph()
p4_3_3.text = "3. Agent uses Search Tool: Queries 'population of Tokyo'."
p4_3_3.level = 1
p4_3_4 = tf4.add_paragraph()
p4_3_4.text = "4. Agent synthesizes: Provides the answer based on search results."
p4_3_4.level = 1

# Slide 5: Key Takeaways & Code Concept
slide5 = prs.slides.add_slide(slide_layout_content)
set_slide_background(slide5, white_color)

title5 = slide5.shapes.title
title5.text = "Building an Agent (Conceptual)"

body5 = slide5.shapes.placeholders[1]
tf5 = body5.text_frame
tf5.clear()

p5_1 = tf5.add_paragraph()
p5_1.text = "Conceptual Langchain Code Snippet:"
p5_1.level = 0

# Add a text box for the code snippet
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(3)
textbox = slide5.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

code_text = """
from langchain.llms import OpenAI
from langchain.agents import initialize_agent, Tool
from langchain.agents import AgentType

# 1. Define your LLM
llm = OpenAI(temperature=0)

# 2. Define your tools
tools = [
    Tool(
        name="Search",
        func=lambda query: "Search results for " + query, # Placeholder for actual search
        description="useful for when you need to answer questions about current events or facts"
    )
]

# 3. Initialize the agent
agent = initialize_agent(tools, llm, agent=AgentType.ZERO_SHOT_REACT_DESCRIPTION, verbose=True)

# 4. Run the agent with a query
# agent.run("What is the capital of France and its current population?")
"""
p_code = text_frame.add_paragraph()
p_code.text = code_text
p_code.font.name = 'Consolas' # A common monospace font for code
p_code.font.size = prs.slide_height * 0.02 # Adjust font size dynamically

p5_2 = tf5.add_paragraph()
p5_2.text = "Agentic AI empowers LLMs to go beyond simple text generation."
p5_2.level = 0

p5_3 = tf5.add_paragraph()
p5_3.text = "Langchain simplifies the creation of such intelligent systems."
p5_3.level = 0

# Save the presentation
prs.save("output.pptx")